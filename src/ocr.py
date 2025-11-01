"""
OCR and Text Extraction Module

This module handles text extraction from various file types including:
- PDFs (text-based and scanned images)
- Images (JPEG, PNG)
- Excel files (XLSX)
- PowerPoint files (PPTX)
"""

import os
import logging
from pathlib import Path
from typing import Dict, List, Optional, Union
import zipfile

# OCR libraries
import pytesseract
import easyocr
from PIL import Image

# PDF processing
import pdfplumber
from PyPDF2 import PdfReader

# Office document processing
import openpyxl
from pptx import Presentation

# Computer vision
import cv2
import numpy as np

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class TextExtractor:
    """Main class for extracting text from various file formats."""
    
    def __init__(self, ocr_engine: str = "pytesseract"):
        """
        Initialize the TextExtractor.
        
        Args:
            ocr_engine: OCR engine to use ("pytesseract" or "easyocr")
        """
        self.ocr_engine = ocr_engine
        if ocr_engine == "easyocr":
            self.reader = easyocr.Reader(['en'])
        
        # Supported file extensions
        self.supported_extensions = {
            '.pdf', '.jpg', '.jpeg', '.png', '.xlsx', '.pptx', '.zip', '.txt'
        }
    
    def extract_text_from_file(self, file_path: Union[str, Path]) -> Dict[str, str]:
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {extension}")
        
        logger.info(f"Extracting text from {file_path}")
        
        try:
            if extension == '.pdf':
                return self._extract_from_pdf(file_path)
            elif extension in ['.jpg', '.jpeg', '.png']:
                return self._extract_from_image(file_path)
            elif extension == '.xlsx':
                return self._extract_from_excel(file_path)
            elif extension == '.pptx':
                return self._extract_from_powerpoint(file_path)
            elif extension == '.zip':
                return self._extract_from_zip(file_path)
            elif extension == '.txt':
                return self._extract_from_text(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return {
                'text': '',
                'file_type': extension,
                'error': str(e),
                'success': False
            }
    
    def _extract_from_pdf(self, file_path: Path) -> Dict[str, str]:
        """Extract text from PDF files (text-based or scanned)."""
        text_content = ""
        is_scanned = False
        
        try:
            # First, try to extract text directly
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
            
            # If no text found, treat as scanned PDF
            if not text_content.strip():
                is_scanned = True
                text_content = self._extract_from_scanned_pdf(file_path)
        
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyPDF2: {str(e)}")
            # Fallback to PyPDF2
            try:
                reader = PdfReader(file_path)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
                
                if not text_content.strip():
                    is_scanned = True
                    text_content = self._extract_from_scanned_pdf(file_path)
            
            except Exception as e2:
                logger.error(f"Both PDF extraction methods failed: {str(e2)}")
                return {
                    'text': '',
                    'file_type': 'pdf',
                    'error': f"PDF extraction failed: {str(e2)}",
                    'success': False
                }
        
        return {
            'text': text_content,
            'file_type': 'pdf',
            'is_scanned': is_scanned,
            'success': True
        }
    
    def _extract_from_scanned_pdf(self, file_path: Path) -> str:
        """Extract text from scanned PDF using OCR."""
        text_content = ""
        
        try:
            import fitz  # PyMuPDF for converting PDF pages to images
            import io
            
            # Convert PDF pages to images and apply OCR
            pdf_document = fitz.open(file_path)
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                # Convert page to image
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                
                # Convert to PIL Image
                image = Image.open(io.BytesIO(img_data))
                
                # Apply OCR
                page_text = self._apply_ocr(image)
                text_content += page_text + "\n"
            
            pdf_document.close()
        
        except ImportError:
            logger.warning("PyMuPDF not available, using alternative method")
            # Alternative: convert PDF to images using other methods
            text_content = "OCR extraction requires PyMuPDF (fitz) library"
        
        return text_content
    
    def _extract_from_image(self, file_path: Path) -> Dict[str, str]:
        """Extract text from image files using OCR."""
        try:
            # Load image
            image = Image.open(file_path)
            
            # Apply OCR
            text_content = self._apply_ocr(image)
            
            return {
                'text': text_content,
                'file_type': 'image',
                'success': True
            }
        
        except Exception as e:
            return {
                'text': '',
                'file_type': 'image',
                'error': str(e),
                'success': False
            }
    
    def _apply_ocr(self, image: Image.Image) -> str:
        """Apply OCR to an image using the configured engine."""
        try:
            if self.ocr_engine == "pytesseract":
                # Preprocess image for better OCR results
                cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
                gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
                
                # Apply image preprocessing
                processed_image = self._preprocess_image(gray)
                
                # Convert back to PIL Image
                pil_image = Image.fromarray(processed_image)
                
                # Apply OCR with custom config
                custom_config = r'--oem 3 --psm 6'
                text = pytesseract.image_to_string(pil_image, config=custom_config)
                
            elif self.ocr_engine == "easyocr":
                # Convert PIL to numpy array
                cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
                
                # Apply EasyOCR
                results = self.reader.readtext(cv_image)
                text = " ".join([result[1] for result in results])
            
            else:
                raise ValueError(f"Unsupported OCR engine: {self.ocr_engine}")
            
            return text
        
        except Exception as e:
            logger.error(f"OCR failed: {str(e)}")
            return ""
    
    def _preprocess_image(self, image: np.ndarray) -> np.ndarray:
        """Preprocess image for better OCR results."""
        # Apply Gaussian blur to reduce noise
        blurred = cv2.GaussianBlur(image, (5, 5), 0)
        
        # Apply threshold to get binary image
        _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Apply morphological operations to clean up
        kernel = np.ones((1, 1), np.uint8)
        processed = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
        processed = cv2.morphologyEx(processed, cv2.MORPH_OPEN, kernel)
        
        return processed
    
    def _extract_from_excel(self, file_path: Path) -> Dict[str, str]:
        """Extract text from Excel files."""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    
                    if row_text:
                        text_content += " | ".join(row_text) + "\n"
            
            workbook.close()
            
            return {
                'text': text_content,
                'file_type': 'excel',
                'success': True
            }
        
        except Exception as e:
            return {
                'text': '',
                'file_type': 'excel',
                'error': str(e),
                'success': False
            }
    
    def _extract_from_powerpoint(self, file_path: Path) -> Dict[str, str]:
        """Extract text from PowerPoint files."""
        try:
            presentation = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            
            return {
                'text': text_content,
                'file_type': 'powerpoint',
                'success': True
            }
        
        except Exception as e:
            return {
                'text': '',
                'file_type': 'powerpoint',
                'error': str(e),
                'success': False
            }
    
    def _extract_from_zip(self, file_path: Path) -> Dict[str, str]:
        """Extract text from files within a ZIP archive."""
        results = {}
        
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # Extract to temporary directory
                temp_dir = file_path.parent / f"temp_{file_path.stem}"
                temp_dir.mkdir(exist_ok=True)
                
                zip_ref.extractall(temp_dir)
                
                # Process each extracted file
                for extracted_file in temp_dir.rglob('*'):
                    if extracted_file.is_file():
                        file_extension = extracted_file.suffix.lower()
                        
                        if file_extension in self.supported_extensions and file_extension != '.zip':
                            try:
                                result = self.extract_text_from_file(extracted_file)
                                results[str(extracted_file.relative_to(temp_dir))] = result
                            except Exception as e:
                                logger.warning(f"Failed to process {extracted_file}: {str(e)}")
                
                # Clean up temporary directory
                import shutil
                shutil.rmtree(temp_dir)
            
            # Combine all extracted text
            combined_text = ""
            for file_name, result in results.items():
                if result.get('success', False):
                    combined_text += f"\n--- File: {file_name} ---\n"
                    combined_text += result.get('text', '') + "\n"
            
            return {
                'text': combined_text,
                'file_type': 'zip',
                'extracted_files': list(results.keys()),
                'success': True
            }
        
        except Exception as e:
            return {
                'text': '',
                'file_type': 'zip',
                'error': str(e),
                'success': False
            }
    
    def _extract_from_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            return {
                'text': text_content,
                'file_type': 'text',
                'success': True
            }
        
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text_content = f.read()
                
                return {
                    'text': text_content,
                    'file_type': 'text',
                    'success': True
                }
            except Exception as e:
                return {
                    'text': '',
                    'file_type': 'text',
                    'error': str(e),
                    'success': False
                }
        
        except Exception as e:
            return {
                'text': '',
                'file_type': 'text',
                'error': str(e),
                'success': False
            }


def extract_text_from_file(file_path: Union[str, Path], ocr_engine: str = "pytesseract") -> Dict[str, str]:
    """
    Convenience function to extract text from a file.
    
    Args:
        file_path: Path to the file
        ocr_engine: OCR engine to use ("pytesseract" or "easyocr")
        
    Returns:
        Dictionary containing extracted text and metadata
    """
    extractor = TextExtractor(ocr_engine=ocr_engine)
    return extractor.extract_text_from_file(file_path)


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        result = extract_text_from_file(file_path)
        
        if result.get('success', False):
            print("Extracted Text:")
            print("-" * 50)
            print(result['text'])
        else:
            print(f"Error: {result.get('error', 'Unknown error')}")
    else:
        print("Usage: python ocr.py <file_path>")
