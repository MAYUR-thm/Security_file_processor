"""
Report Generator Module

This module generates various types of reports from security analysis results:
- JSON reports with structured data
- CSV reports for spreadsheet analysis
- PowerPoint presentations with summary tables
- Workflow diagrams showing the processing pipeline
"""

import json
import csv
import logging
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np

# Try to import graphviz for workflow diagrams
try:
    import graphviz
    GRAPHVIZ_AVAILABLE = True
except ImportError:
    GRAPHVIZ_AVAILABLE = False
    logging.warning("Graphviz not available. Workflow diagrams will use matplotlib instead.")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportGenerator:
    """Main class for generating various types of reports."""
    
    def __init__(self, output_dir: Union[str, Path] = "output"):
        """
        Initialize the ReportGenerator.
        
        Args:
            output_dir: Directory to save generated reports
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Report templates and styling
        self.report_styles = {
            'colors': {
                'primary': '#2E86AB',
                'secondary': '#A23B72',
                'success': '#F18F01',
                'warning': '#C73E1D',
                'info': '#6C757D',
                'light': '#F8F9FA',
                'dark': '#343A40'
            },
            'fonts': {
                'title': 'Calibri',
                'body': 'Arial',
                'code': 'Courier New'
            }
        }
    
    def generate_comprehensive_report(self, results: List[Dict[str, Any]], 
                                   report_name: str = "security_analysis_report") -> Dict[str, str]:
        """
        Generate a comprehensive report in multiple formats.
        
        Args:
            results: List of analysis results for each processed file
            report_name: Base name for the report files
            
        Returns:
            Dictionary containing paths to generated report files
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = f"{report_name}_{timestamp}"
        
        generated_files = {}
        
        try:
            # Generate JSON report
            json_path = self.generate_json_report(results, f"{base_name}.json")
            generated_files['json'] = json_path
            
            # Generate CSV report
            csv_path = self.generate_csv_report(results, f"{base_name}.csv")
            generated_files['csv'] = csv_path
            
            # Generate PowerPoint report
            pptx_path = self.generate_powerpoint_report(results, f"{base_name}.pptx")
            generated_files['powerpoint'] = pptx_path
            
            # Generate workflow diagram
            diagram_path = self.generate_workflow_diagram(f"{base_name}_workflow.png")
            generated_files['workflow_diagram'] = diagram_path
            
            # Generate summary statistics
            stats_path = self.generate_statistics_report(results, f"{base_name}_statistics.json")
            generated_files['statistics'] = stats_path
            
            logger.info(f"Comprehensive report generated: {len(generated_files)} files created")
            
        except Exception as e:
            logger.error(f"Error generating comprehensive report: {str(e)}")
            raise
        
        return generated_files
    
    def generate_json_report(self, results: List[Dict[str, Any]], 
                           filename: str = "security_report.json") -> str:
        """
        Generate a structured JSON report.
        
        Args:
            results: Analysis results
            filename: Output filename
            
        Returns:
            Path to generated JSON file
        """
        output_path = self.output_dir / filename
        
        # Create comprehensive JSON structure
        report_data = {
            'metadata': {
                'generated_at': datetime.now().isoformat(),
                'generator': 'Security File Processor',
                'version': '1.0.0',
                'total_files_processed': len(results)
            },
            'summary': self._generate_summary_statistics(results),
            'files': results,
            'aggregated_findings': self._aggregate_findings(results),
            'recommendations': self._aggregate_recommendations(results)
        }
        
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(report_data, f, indent=2, ensure_ascii=False, default=str)
            
            logger.info(f"JSON report saved to: {output_path}")
            return str(output_path)
        
        except Exception as e:
            logger.error(f"Error generating JSON report: {str(e)}")
            raise
    
    def generate_csv_report(self, results: List[Dict[str, Any]], 
                          filename: str = "security_report.csv") -> str:
        """
        Generate a CSV report suitable for spreadsheet analysis.
        
        Args:
            results: Analysis results
            filename: Output filename
            
        Returns:
            Path to generated CSV file
        """
        output_path = self.output_dir / filename
        
        # Flatten results for CSV format
        csv_data = []
        
        for result in results:
            base_row = {
                'file_name': result.get('file_name', ''),
                'file_type': result.get('file_type', ''),
                'file_description': result.get('file_description', ''),
                'processing_timestamp': result.get('processing_timestamp', ''),
                'text_length': result.get('text_length', 0),
                'pii_risk_level': result.get('pii_analysis', {}).get('risk_level', ''),
                'pii_items_found': result.get('pii_analysis', {}).get('pii_count', 0),
                'security_confidence': result.get('security_analysis', {}).get('summary', {}).get('confidence_score', 0),
                'total_security_findings': result.get('security_analysis', {}).get('summary', {}).get('total_findings', 0)
            }
            
            # Add key findings as a concatenated string
            key_findings = result.get('security_analysis', {}).get('key_findings', [])
            base_row['key_findings'] = ' | '.join(key_findings[:5])  # Limit to first 5
            
            # Add security categories found
            categories_found = result.get('security_analysis', {}).get('summary', {}).get('categories_found', [])
            base_row['security_categories'] = ', '.join(categories_found)
            
            csv_data.append(base_row)
        
        try:
            # Write CSV file
            if csv_data:
                df = pd.DataFrame(csv_data)
                df.to_csv(output_path, index=False, encoding='utf-8')
            else:
                # Create empty CSV with headers
                headers = ['file_name', 'file_type', 'file_description', 'processing_timestamp',
                          'text_length', 'pii_risk_level', 'pii_items_found', 'security_confidence',
                          'total_security_findings', 'key_findings', 'security_categories']
                
                with open(output_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerow(headers)
            
            logger.info(f"CSV report saved to: {output_path}")
            return str(output_path)
        
        except Exception as e:
            logger.error(f"Error generating CSV report: {str(e)}")
            raise
    
    def generate_powerpoint_report(self, results: List[Dict[str, Any]], 
                                 filename: str = "security_report.pptx") -> str:
        """
        Generate a PowerPoint presentation with analysis results.
        
        Args:
            results: Analysis results
            filename: Output filename
            
        Returns:
            Path to generated PowerPoint file
        """
        output_path = self.output_dir / filename
        
        try:
            # Create presentation
            prs = Presentation()
            
            # Slide 1: Title slide
            self._add_title_slide(prs, results)
            
            # Slide 2: Executive summary
            self._add_executive_summary_slide(prs, results)
            
            # Slide 3: Processing statistics
            self._add_statistics_slide(prs, results)
            
            # Slide 4+: Detailed results table
            self._add_detailed_results_slides(prs, results)
            
            # Slide N: Recommendations
            self._add_recommendations_slide(prs, results)
            
            # Save presentation
            prs.save(output_path)
            
            logger.info(f"PowerPoint report saved to: {output_path}")
            return str(output_path)
        
        except Exception as e:
            logger.error(f"Error generating PowerPoint report: {str(e)}")
            raise
    
    def _add_title_slide(self, prs: Presentation, results: List[Dict[str, Any]]) -> None:
        """Add title slide to PowerPoint presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Security File Processing Report"
        
        subtitle_text = f"""
        Analysis of {len(results)} files
        Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}
        
        Automated Security Content Analysis
        PII Detection and Cleansing Results
        """
        
        subtitle.text = subtitle_text.strip()
    
    def _add_executive_summary_slide(self, prs: Presentation, results: List[Dict[str, Any]]) -> None:
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Calculate summary statistics
        summary_stats = self._generate_summary_statistics(results)
        
        content = slide.placeholders[1]
        content_text = f"""
        Files Processed: {summary_stats['total_files']}
        
        PII Analysis:
        • High Risk Files: {summary_stats['pii_stats']['high_risk_files']}
        • Total PII Items Found: {summary_stats['pii_stats']['total_pii_items']}
        
        Security Analysis:
        • Files with Security Content: {summary_stats['security_stats']['files_with_security_content']}
        • Total Security Findings: {summary_stats['security_stats']['total_findings']}
        • Average Confidence Score: {summary_stats['security_stats']['avg_confidence']:.2f}
        
        File Types Processed:
        {self._format_file_types(summary_stats['file_types'])}
        """
        
        content.text = content_text.strip()
    
    def _add_statistics_slide(self, prs: Presentation, results: List[Dict[str, Any]]) -> None:
        """Add processing statistics slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Processing Statistics"
        
        # Generate statistics chart and add to slide
        chart_path = self._create_statistics_chart(results)
        
        if chart_path and Path(chart_path).exists():
            # Add chart image to slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            
            slide.shapes.add_picture(chart_path, left, top, width, height)
        else:
            # Fallback to text-based statistics
            content = slide.placeholders[1]
            stats = self._generate_summary_statistics(results)
            
            content_text = f"""
            Processing Overview:
            • Total Files: {stats['total_files']}
            • Successful Processing: {stats['successful_processing']}
            • Files with Errors: {stats['files_with_errors']}
            
            Security Categories Found:
            {self._format_security_categories(stats['security_stats']['categories_distribution'])}
            
            PII Risk Distribution:
            {self._format_pii_distribution(stats['pii_stats']['risk_distribution'])}
            """
            
            content.text = content_text.strip()
    
    def _add_detailed_results_slides(self, prs: Presentation, results: List[Dict[str, Any]]) -> None:
        """Add detailed results table slides."""
        # Create table data
        table_data = []
        headers = ['File Name', 'Type', 'Description', 'Key Findings']
        
        for result in results:
            key_findings = result.get('security_analysis', {}).get('key_findings', [])
            findings_text = '; '.join(key_findings[:3]) if key_findings else 'None'
            
            row = [
                result.get('file_name', 'Unknown'),
                result.get('file_type', 'Unknown'),
                result.get('file_description', 'No description'),
                findings_text
            ]
            table_data.append(row)
        
        # Split data into chunks for multiple slides if needed
        rows_per_slide = 8
        
        for i in range(0, len(table_data), rows_per_slide):
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            title_frame = title_shape.text_frame
            title_frame.text = f"Detailed Results (Page {i//rows_per_slide + 1})"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Add table
            chunk = table_data[i:i + rows_per_slide]
            
            rows = len(chunk) + 1  # +1 for header
            cols = len(headers)
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            table.columns[0].width = Inches(2)  # File Name
            table.columns[1].width = Inches(1)  # Type
            table.columns[2].width = Inches(2.5)  # Description
            table.columns[3].width = Inches(3.5)  # Key Findings
            
            # Add headers
            for j, header in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(47, 117, 181)  # Blue background
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
            
            # Add data rows
            for row_idx, row_data in enumerate(chunk, 1):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_data)[:100]  # Limit cell text length
                    
                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
    
    def _add_recommendations_slide(self, prs: Presentation, results: List[Dict[str, Any]]) -> None:
        """Add recommendations slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Recommendations"
        
        # Aggregate recommendations
        all_recommendations = self._aggregate_recommendations(results)
        
        content = slide.placeholders[1]
        
        if all_recommendations:
            recommendations_text = "\n".join([f"• {rec}" for rec in all_recommendations[:10]])
        else:
            recommendations_text = "• No specific recommendations generated\n• Consider implementing regular security assessments\n• Maintain up-to-date security policies"
        
        content.text = recommendations_text
    
    def generate_workflow_diagram(self, filename: str = "workflow_diagram.png") -> str:
        """
        Generate a workflow diagram showing the processing pipeline.
        
        Args:
            filename: Output filename
            
        Returns:
            Path to generated diagram file
        """
        output_path = self.output_dir / filename
        
        try:
            if GRAPHVIZ_AVAILABLE:
                return self._generate_graphviz_diagram(output_path)
            else:
                return self._generate_matplotlib_diagram(output_path)
        
        except Exception as e:
            logger.error(f"Error generating workflow diagram: {str(e)}")
            # Create a simple text-based diagram as fallback
            return self._generate_text_diagram(output_path)
    
    def _generate_graphviz_diagram(self, output_path: Path) -> str:
        """Generate workflow diagram using Graphviz."""
        dot = graphviz.Digraph(comment='Security File Processing Workflow')
        dot.attr(rankdir='TB', size='10,8')
        
        # Define nodes
        nodes = [
            ('input', 'File Input\n(PDF, Images, Excel, PowerPoint)', 'lightblue'),
            ('ocr', 'Text Extraction\n(OCR, Document Parsing)', 'lightgreen'),
            ('pii', 'PII Detection\n& Cleansing', 'yellow'),
            ('logo', 'Logo Detection\n& Masking', 'orange'),
            ('analysis', 'Security Content\nAnalysis', 'pink'),
            ('report', 'Report Generation\n(JSON, CSV, PowerPoint)', 'lightcoral')
        ]
        
        for node_id, label, color in nodes:
            dot.node(node_id, label, style='filled', fillcolor=color, shape='box')
        
        # Define edges
        edges = [
            ('input', 'ocr'),
            ('ocr', 'pii'),
            ('input', 'logo'),
            ('pii', 'analysis'),
            ('analysis', 'report')
        ]
        
        for src, dst in edges:
            dot.edge(src, dst)
        
        # Render diagram
        dot.render(str(output_path.with_suffix('')), format='png', cleanup=True)
        
        logger.info(f"Graphviz workflow diagram saved to: {output_path}")
        return str(output_path)
    
    def _generate_matplotlib_diagram(self, output_path: Path) -> str:
        """Generate workflow diagram using matplotlib."""
        fig, ax = plt.subplots(1, 1, figsize=(12, 8))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 8)
        ax.axis('off')
        
        # Define boxes
        boxes = [
            (5, 7, 'File Input\n(Multiple Formats)', 'lightblue'),
            (2, 5.5, 'OCR & Text\nExtraction', 'lightgreen'),
            (8, 5.5, 'Logo Detection\n& Masking', 'orange'),
            (5, 4, 'PII Detection\n& Cleansing', 'yellow'),
            (5, 2.5, 'Security Content\nAnalysis', 'pink'),
            (5, 1, 'Report Generation', 'lightcoral')
        ]
        
        # Draw boxes
        for x, y, text, color in boxes:
            box = FancyBboxPatch((x-0.8, y-0.4), 1.6, 0.8, 
                               boxstyle="round,pad=0.1", 
                               facecolor=color, 
                               edgecolor='black',
                               linewidth=1)
            ax.add_patch(box)
            ax.text(x, y, text, ha='center', va='center', fontsize=10, weight='bold')
        
        # Draw arrows
        arrows = [
            (5, 6.6, 2, 5.9),  # Input to OCR
            (5, 6.6, 8, 5.9),  # Input to Logo
            (2, 5.1, 4.2, 4.4),  # OCR to PII
            (5, 3.6, 5, 2.9),  # PII to Analysis
            (5, 2.1, 5, 1.4)   # Analysis to Report
        ]
        
        for x1, y1, x2, y2 in arrows:
            ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                       arrowprops=dict(arrowstyle='->', lw=2, color='black'))
        
        plt.title('Security File Processing Workflow', fontsize=16, weight='bold', pad=20)
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        logger.info(f"Matplotlib workflow diagram saved to: {output_path}")
        return str(output_path)
    
    def _generate_text_diagram(self, output_path: Path) -> str:
        """Generate a simple text-based diagram as fallback."""
        diagram_text = """
Security File Processing Workflow
=================================

┌─────────────────────┐
│    File Input       │
│  (Multiple Formats) │
└──────────┬──────────┘
           │
           ├─────────────────────────────┐
           │                             │
           ▼                             ▼
┌─────────────────────┐         ┌─────────────────────┐
│  OCR & Text         │         │  Logo Detection     │
│  Extraction         │         │  & Masking          │
└──────────┬──────────┘         └─────────────────────┘
           │
           ▼
┌─────────────────────┐
│  PII Detection      │
│  & Cleansing        │
└──────────┬──────────┘
           │
           ▼
┌─────────────────────┐
│  Security Content   │
│  Analysis           │
└──────────┬──────────┘
           │
           ▼
┌─────────────────────┐
│  Report Generation  │
│  (JSON, CSV, PPTX)  │
└─────────────────────┘
        """
        
        # Save as text file instead of image
        text_output_path = output_path.with_suffix('.txt')
        with open(text_output_path, 'w', encoding='utf-8') as f:
            f.write(diagram_text.strip())
        
        logger.info(f"Text workflow diagram saved to: {text_output_path}")
        return str(text_output_path)
    
    def generate_statistics_report(self, results: List[Dict[str, Any]], 
                                 filename: str = "statistics.json") -> str:
        """Generate detailed statistics report."""
        output_path = self.output_dir / filename
        
        stats = self._generate_summary_statistics(results)
        
        # Add more detailed statistics
        detailed_stats = {
            **stats,
            'detailed_analysis': {
                'processing_times': self._calculate_processing_times(results),
                'file_size_analysis': self._analyze_file_sizes(results),
                'security_category_details': self._analyze_security_categories(results),
                'pii_type_distribution': self._analyze_pii_types(results)
            }
        }
        
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(detailed_stats, f, indent=2, ensure_ascii=False, default=str)
            
            logger.info(f"Statistics report saved to: {output_path}")
            return str(output_path)
        
        except Exception as e:
            logger.error(f"Error generating statistics report: {str(e)}")
            raise
    
    def _generate_summary_statistics(self, results: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Generate summary statistics from results."""
        if not results:
            return {
                'total_files': 0,
                'successful_processing': 0,
                'files_with_errors': 0,
                'file_types': {},
                'pii_stats': {},
                'security_stats': {}
            }
        
        # Basic counts
        total_files = len(results)
        successful = sum(1 for r in results if r.get('success', True))
        errors = total_files - successful
        
        # File types
        file_types = {}
        for result in results:
            file_type = result.get('file_type', 'unknown')
            file_types[file_type] = file_types.get(file_type, 0) + 1
        
        # PII statistics
        pii_stats = {
            'high_risk_files': 0,
            'total_pii_items': 0,
            'risk_distribution': {'LOW': 0, 'MEDIUM': 0, 'HIGH': 0, 'CRITICAL': 0}
        }
        
        for result in results:
            pii_analysis = result.get('pii_analysis', {})
            risk_level = pii_analysis.get('risk_level', 'LOW')
            pii_count = pii_analysis.get('pii_count', 0)
            
            if risk_level in ['HIGH', 'CRITICAL']:
                pii_stats['high_risk_files'] += 1
            
            pii_stats['total_pii_items'] += pii_count
            pii_stats['risk_distribution'][risk_level] = pii_stats['risk_distribution'].get(risk_level, 0) + 1
        
        # Security statistics
        security_stats = {
            'files_with_security_content': 0,
            'total_findings': 0,
            'avg_confidence': 0.0,
            'categories_distribution': {}
        }
        
        confidence_scores = []
        
        for result in results:
            security_analysis = result.get('security_analysis', {})
            summary = security_analysis.get('summary', {})
            
            findings = summary.get('total_findings', 0)
            confidence = summary.get('confidence_score', 0)
            categories = summary.get('categories_found', [])
            
            if findings > 0:
                security_stats['files_with_security_content'] += 1
            
            security_stats['total_findings'] += findings
            confidence_scores.append(confidence)
            
            for category in categories:
                security_stats['categories_distribution'][category] = security_stats['categories_distribution'].get(category, 0) + 1
        
        if confidence_scores:
            security_stats['avg_confidence'] = sum(confidence_scores) / len(confidence_scores)
        
        return {
            'total_files': total_files,
            'successful_processing': successful,
            'files_with_errors': errors,
            'file_types': file_types,
            'pii_stats': pii_stats,
            'security_stats': security_stats
        }
    
    def _aggregate_findings(self, results: List[Dict[str, Any]]) -> Dict[str, List[str]]:
        """Aggregate all findings across files."""
        aggregated = {
            'all_key_findings': [],
            'security_categories': [],
            'pii_types': [],
            'high_confidence_findings': []
        }
        
        for result in results:
            # Security findings
            security_analysis = result.get('security_analysis', {})
            key_findings = security_analysis.get('key_findings', [])
            categories = security_analysis.get('summary', {}).get('categories_found', [])
            confidence = security_analysis.get('summary', {}).get('confidence_score', 0)
            
            aggregated['all_key_findings'].extend(key_findings)
            aggregated['security_categories'].extend(categories)
            
            if confidence > 0.7:
                aggregated['high_confidence_findings'].extend(key_findings)
            
            # PII types
            pii_analysis = result.get('pii_analysis', {})
            pii_findings = pii_analysis.get('pii_findings', {})
            
            if 'statistics' in pii_findings:
                entity_types = pii_findings['statistics'].get('entity_types', {})
                pattern_types = pii_findings['statistics'].get('pattern_types', {})
                
                aggregated['pii_types'].extend(list(entity_types.keys()))
                aggregated['pii_types'].extend(list(pattern_types.keys()))
        
        # Remove duplicates and sort
        for key in aggregated:
            aggregated[key] = sorted(list(set(aggregated[key])))
        
        return aggregated
    
    def _aggregate_recommendations(self, results: List[Dict[str, Any]]) -> List[str]:
        """Aggregate recommendations from all analyses."""
        all_recommendations = []
        
        for result in results:
            # Security recommendations
            security_analysis = result.get('security_analysis', {})
            security_recs = security_analysis.get('recommendations', [])
            all_recommendations.extend(security_recs)
            
            # PII recommendations
            pii_analysis = result.get('pii_analysis', {})
            pii_recs = pii_analysis.get('recommendations', [])
            all_recommendations.extend(pii_recs)
        
        # Remove duplicates and return top recommendations
        unique_recommendations = list(set(all_recommendations))
        return unique_recommendations[:15]  # Top 15 recommendations
    
    def _create_statistics_chart(self, results: List[Dict[str, Any]]) -> Optional[str]:
        """Create a statistics chart and return the file path."""
        try:
            stats = self._generate_summary_statistics(results)
            
            fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(12, 8))
            
            # Chart 1: File types
            if stats['file_types']:
                ax1.pie(stats['file_types'].values(), labels=stats['file_types'].keys(), autopct='%1.1f%%')
                ax1.set_title('File Types Distribution')
            
            # Chart 2: PII risk distribution
            if stats['pii_stats']['risk_distribution']:
                risk_data = stats['pii_stats']['risk_distribution']
                ax2.bar(risk_data.keys(), risk_data.values(), color=['green', 'yellow', 'orange', 'red'])
                ax2.set_title('PII Risk Distribution')
                ax2.set_ylabel('Number of Files')
            
            # Chart 3: Security categories
            if stats['security_stats']['categories_distribution']:
                cat_data = stats['security_stats']['categories_distribution']
                ax3.barh(list(cat_data.keys()), list(cat_data.values()))
                ax3.set_title('Security Categories Found')
                ax3.set_xlabel('Number of Files')
            
            # Chart 4: Processing summary
            processing_data = [stats['successful_processing'], stats['files_with_errors']]
            ax4.pie(processing_data, labels=['Successful', 'Errors'], autopct='%1.1f%%', colors=['lightgreen', 'lightcoral'])
            ax4.set_title('Processing Results')
            
            plt.tight_layout()
            
            chart_path = self.output_dir / 'statistics_chart.png'
            plt.savefig(chart_path, dpi=300, bbox_inches='tight')
            plt.close()
            
            return str(chart_path)
        
        except Exception as e:
            logger.warning(f"Could not create statistics chart: {str(e)}")
            return None
    
    def _format_file_types(self, file_types: Dict[str, int]) -> str:
        """Format file types for display."""
        if not file_types:
            return "No files processed"
        
        return "\n".join([f"• {file_type}: {count}" for file_type, count in file_types.items()])
    
    def _format_security_categories(self, categories: Dict[str, int]) -> str:
        """Format security categories for display."""
        if not categories:
            return "No security categories found"
        
        return "\n".join([f"• {cat.replace('_', ' ').title()}: {count}" for cat, count in categories.items()])
    
    def _format_pii_distribution(self, distribution: Dict[str, int]) -> str:
        """Format PII distribution for display."""
        if not distribution:
            return "No PII analysis data"
        
        return "\n".join([f"• {risk}: {count}" for risk, count in distribution.items()])
    
    def _calculate_processing_times(self, results: List[Dict[str, Any]]) -> Dict[str, float]:
        """Calculate processing time statistics."""
        # This would require timing data from the actual processing
        # For now, return placeholder data
        return {
            'avg_processing_time': 0.0,
            'total_processing_time': 0.0,
            'fastest_file': 0.0,
            'slowest_file': 0.0
        }
    
    def _analyze_file_sizes(self, results: List[Dict[str, Any]]) -> Dict[str, int]:
        """Analyze file size distribution."""
        # This would require file size data
        # For now, return placeholder data based on text length
        sizes = {'small': 0, 'medium': 0, 'large': 0}
        
        for result in results:
            text_length = result.get('text_length', 0)
            if text_length < 1000:
                sizes['small'] += 1
            elif text_length < 10000:
                sizes['medium'] += 1
            else:
                sizes['large'] += 1
        
        return sizes
    
    def _analyze_security_categories(self, results: List[Dict[str, Any]]) -> Dict[str, Dict[str, float]]:
        """Analyze security categories in detail."""
        category_analysis = {}
        
        for result in results:
            security_analysis = result.get('security_analysis', {})
            categories = security_analysis.get('categories', {})
            
            for category, data in categories.items():
                if category not in category_analysis:
                    category_analysis[category] = {
                        'total_files': 0,
                        'avg_confidence': 0.0,
                        'total_findings': 0
                    }
                
                category_analysis[category]['total_files'] += 1
                category_analysis[category]['avg_confidence'] += data.get('confidence', 0)
                category_analysis[category]['total_findings'] += data.get('finding_count', 0)
        
        # Calculate averages
        for category in category_analysis:
            total_files = category_analysis[category]['total_files']
            if total_files > 0:
                category_analysis[category]['avg_confidence'] /= total_files
        
        return category_analysis
    
    def _analyze_pii_types(self, results: List[Dict[str, Any]]) -> Dict[str, int]:
        """Analyze PII types distribution."""
        pii_types = {}
        
        for result in results:
            pii_analysis = result.get('pii_analysis', {})
            pii_findings = pii_analysis.get('pii_findings', {})
            
            if 'statistics' in pii_findings:
                entity_types = pii_findings['statistics'].get('entity_types', {})
                pattern_types = pii_findings['statistics'].get('pattern_types', {})
                
                for pii_type, count in entity_types.items():
                    pii_types[pii_type] = pii_types.get(pii_type, 0) + count
                
                for pii_type, count in pattern_types.items():
                    pii_types[pii_type] = pii_types.get(pii_type, 0) + count
        
        return pii_types


def generate_report(results: List[Dict[str, Any]], 
                   output_dir: Union[str, Path] = "output",
                   report_name: str = "security_analysis_report") -> Dict[str, str]:
    """
    Convenience function to generate comprehensive reports.
    
    Args:
        results: Analysis results
        output_dir: Output directory
        report_name: Base name for report files
        
    Returns:
        Dictionary containing paths to generated files
    """
    generator = ReportGenerator(output_dir)
    return generator.generate_comprehensive_report(results, report_name)


if __name__ == "__main__":
    # Example usage with sample data
    sample_results = [
        {
            'file_name': 'security_policy.pdf',
            'file_type': 'pdf',
            'file_description': 'Contains IAM policy configurations',
            'processing_timestamp': datetime.now().isoformat(),
            'text_length': 5000,
            'success': True,
            'pii_analysis': {
                'risk_level': 'MEDIUM',
                'pii_count': 3,
                'recommendations': ['Review and mask identified PII']
            },
            'security_analysis': {
                'summary': {
                    'total_findings': 15,
                    'confidence_score': 0.85,
                    'categories_found': ['iam_policies', 'firewall_rules']
                },
                'key_findings': [
                    'IAM Policy: Allow user access to S3 bucket',
                    'Firewall Rule: Deny TCP port 22'
                ],
                'recommendations': ['Review IAM policies for compliance']
            }
        }
    ]
    
    print("Generating sample reports...")
    generator = ReportGenerator()
    files = generator.generate_comprehensive_report(sample_results, "sample_report")
    
    print("Generated files:")
    for report_type, file_path in files.items():
        print(f"- {report_type}: {file_path}")
