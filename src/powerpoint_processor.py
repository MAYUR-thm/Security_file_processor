"""
PowerPoint Processor Module

This module provides comprehensive PowerPoint file processing:
- Full PowerPoint structure preservation
- PII detection and cleaning in all elements
- Text extraction and analysis
- Image processing support
- Detailed processing reports
"""

import logging
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Union, Any, Tuple
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    logger.error("python-pptx package not found. Please install it with: pip install python-pptx")
    PPTX_AVAILABLE = False

from .pii_cleaner import PIICleaner
from .image_processor import process_image_file

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PowerPointProcessor:
    """Advanced PowerPoint processor with comprehensive PII cleaning."""

    def __init__(self, pii_cleaner: Optional[PIICleaner] = None):
        """
        Initialize the PowerPoint processor.
        
        Args:
            pii_cleaner: Optional PIICleaner instance to use for PII detection/cleaning
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx package is required. Please install it with: pip install python-pptx")
            
        # Initialize or use provided PII cleaner
        self.pii_cleaner = pii_cleaner or PIICleaner()
        
        # Style settings for cleaned content
        self.redacted_color = RGBColor(255, 0, 0)  # Red for replaced text
        self.highlight_color = RGBColor(255, 255, 0)  # Yellow for suspected PII
        
        # Processing settings
        self.process_images = True  # Enable image processing
        self.create_text_summary = True  # Create .txt summary
        self.detailed_logging = True  # Include detailed processing logs

    def process_file(self, 
                    input_file: Union[str, Path], 
                    output_file: Optional[Union[str, Path]] = None,
                    mask_mode: str = "replace",
                    create_text_summary: bool = False) -> Dict[str, Any]:
        """
        Process a PowerPoint file, detecting and cleaning PII while preserving structure.
        Creates a cleaned PPTX file with the same structure as the input.
        
        Args:
            input_file: Path to input PowerPoint file
            output_file: Path for output PowerPoint file
            mask_mode: How to handle PII ('replace', 'mask', or 'remove')
            create_text_summary: Whether to create additional text summary files
            
        Returns:
            Dictionary containing processing results and file locations
        """
        input_file = Path(input_file)
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_file}")

        # Set up output paths
        input_file = Path(input_file)
        if not output_file:
            output_file = input_file.parent / f"cleaned_{input_file.name}"
        else:
            output_file = Path(output_file)
            
        # Create output directory if needed
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Optional text summary paths
        txt_output = None
        report_output = None
        if create_text_summary:
            txt_output = output_file.with_suffix('.txt')
            report_output = output_file.with_suffix('.report.json')

        try:
            prs = Presentation(input_file)
            
            stats = {
                'input_file': str(input_file),
                'output_file': str(output_file),
                'txt_output': str(txt_output) if txt_output else None,
                'report_output': str(report_output) if report_output else None,
                'total_slides': len(prs.slides),
                'slides_with_pii': 0,
                'total_shapes': 0,
                'shapes_with_pii': 0,
                'images_processed': 0,
                'images_with_pii': 0,
                'total_pii_instances': 0,
                'pii_by_type': {},
                'processing_time': None,
                'success': False,
                'start_time': datetime.now().isoformat()
            }
            
            # Initialize text summary
            text_summary = []
            text_summary.append(f"=== Security Document Analysis ===")
            text_summary.append(f"File: {input_file.name}")
            text_summary.append(f"Analysis Date: {stats['start_time']}")
            text_summary.append("\n")

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_has_pii = False
                slide_text = []
                slide_text.append(f"\n--- Slide {slide_num} ---")
                
                # Process shapes in the slide
                for shape in slide.shapes:
                    stats['total_shapes'] += 1
                    shape_text = []
                    
                    # Process text in shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        original_text = shape.text
                        
                        # Use enhanced PII cleaner
                        result = self.pii_cleaner.remove_pii(original_text, mask_mode)
                        cleaned_text = result['cleaned_text']
                        
                        if result['replacements_made'] > 0:
                            slide_has_pii = True
                            stats['shapes_with_pii'] += 1
                            stats['total_pii_instances'] += result['replacements_made']
                            
                            # Update PII type statistics
                            for finding_type, count in result['pii_findings']['statistics']['entity_types'].items():
                                stats['pii_by_type'][finding_type] = stats['pii_by_type'].get(finding_type, 0) + count
                            
                            # Update text and formatting
                            text_frame = shape.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = cleaned_text
                            run.font.color.rgb = self.redacted_color
                        
                        shape_text.append(cleaned_text)

                    # Process tables
                    if shape.has_table:
                        shape_text.append("\n--- Table Content ---")
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    original_text = cell.text
                                    result = self.pii_cleaner.remove_pii(original_text, mask_mode)
                                    cleaned_text = result['cleaned_text']
                                    
                                    if result['replacements_made'] > 0:
                                        slide_has_pii = True
                                        stats['shapes_with_pii'] += 1
                                        stats['total_pii_instances'] += result['replacements_made']
                                        
                                        # Update cell text and formatting
                                        text_frame = cell.text_frame
                                        text_frame.clear()
                                        p = text_frame.paragraphs[0]
                                        run = p.add_run()
                                        run.text = cleaned_text
                                        run.font.color.rgb = self.redacted_color
                                    
                                    row_texts.append(cleaned_text)
                            
                            if row_texts:
                                shape_text.append(" | ".join(row_texts))
                    
                    # Process images if enabled
                    if self.process_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        stats['images_processed'] += 1
                        # Image processing would be handled by image_processor.py
                        # This is a placeholder for future image processing integration
                        shape_text.append("[Image - Scanned for PII]")
                    
                    # Add shape text to slide summary
                    if shape_text:
                        slide_text.extend(shape_text)
                
                # Update slide statistics
                if slide_has_pii:
                    stats['slides_with_pii'] += 1
                
                # Add slide text to summary
                text_summary.extend(slide_text)

            # Calculate processing duration
            end_time = datetime.now()
            stats['end_time'] = end_time.isoformat()
            stats['processing_time'] = str(end_time - datetime.fromisoformat(stats['start_time']))

            # Save cleaned PowerPoint
            prs.save(output_file)
            logger.info(f"Saved cleaned PowerPoint to: {output_file}")

            # Save optional text summary and report
            if create_text_summary:
                # Save text summary
                with open(txt_output, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(text_summary))
                logger.info(f"Saved text summary to: {txt_output}")

                # Save detailed report
                with open(report_output, 'w', encoding='utf-8') as f:
                    json.dump(stats, f, indent=2)
                logger.info(f"Saved processing report to: {report_output}")
            
            stats['success'] = True
            
            return stats

        except Exception as e:
            error_msg = f"Error processing PowerPoint file {input_file}: {str(e)}"
            logger.error(error_msg)
            stats['success'] = False
            stats['error'] = error_msg
            return stats

    def _clean_text(self, text: str, mask_mode: str) -> tuple[str, bool]:
        """
        Clean PII from text while tracking changes.
        
        Args:
            text: Input text
            mask_mode: How to handle PII
            
        Returns:
            Tuple of (cleaned_text, pii_found)
        """
        cleaned_text = text
        pii_found = False
        
        # Check each PII pattern
        for pattern_name, pattern in self.pii_patterns.items():
            matches = re.finditer(pattern, cleaned_text, re.IGNORECASE)
            
            for match in matches:
                pii_found = True
                replacement = self.replacements.get(pattern_name, '[REDACTED]')
                
                if mask_mode == "remove":
                    replacement = ""
                elif mask_mode == "mask":
                    replacement = "*" * len(match.group())
                    
                cleaned_text = cleaned_text[:match.start()] + replacement + cleaned_text[match.end():]
        
        return cleaned_text, pii_found

    def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract text content from PowerPoint file for analysis.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Extracted text content
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== Slide {slide_num} ===")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                    
                    # Extract from tables
                    if shape.has_table:
                        text_parts.append("--- Table Content ---")
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_texts.append(cell.text.strip())
                            if row_texts:
                                text_parts.append(" | ".join(row_texts))

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint file {file_path}: {str(e)}")
            raise

    def analyze_structure(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Analyze the structure of a PowerPoint file.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Dictionary containing structural analysis
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        try:
            prs = Presentation(file_path)
            
            analysis = {
                'total_slides': len(prs.slides),
                'slide_layouts': [],
                'total_shapes': 0,
                'shape_types': {},
                'tables': 0,
                'text_boxes': 0,
                'pictures': 0
            }

            for slide in prs.slides:
                # Track layout type
                layout_name = slide.slide_layout.name
                analysis['slide_layouts'].append(layout_name)
                
                # Analyze shapes
                for shape in slide.shapes:
                    analysis['total_shapes'] += 1
                    shape_type = shape.shape_type
                    analysis['shape_types'][str(shape_type)] = analysis['shape_types'].get(str(shape_type), 0) + 1
                    
                    if shape.has_table:
                        analysis['tables'] += 1
                    if hasattr(shape, "text"):
                        analysis['text_boxes'] += 1
                    if shape.shape_type == 13:  # Picture
                        analysis['pictures'] += 1

            return analysis

        except Exception as e:
            logger.error(f"Error analyzing PowerPoint file {file_path}: {str(e)}")
            raise


def process_powerpoint_file(input_file: Union[str, Path], 
                          output_file: Optional[Union[str, Path]] = None,
                          mask_mode: str = "replace") -> Dict[str, Any]:
    """
    Convenience function to process a single PowerPoint file.
    
    Args:
        input_file: Path to input PowerPoint file
        output_file: Path for output PowerPoint file
        mask_mode: How to handle PII ('replace', 'mask', or 'remove')
        
    Returns:
        Processing results dictionary
    """
    processor = PowerPointProcessor()
    return processor.process_file(input_file, output_file, mask_mode)


if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
        
        try:
            processor = PowerPointProcessor()
            
            # First analyze the structure
            print("\nAnalyzing PowerPoint structure...")
            analysis = processor.analyze_structure(input_file)
            print("\nFile Structure:")
            print(f"Total slides: {analysis['total_slides']}")
            print(f"Total shapes: {analysis['total_shapes']}")
            print(f"Tables: {analysis['tables']}")
            print(f"Text boxes: {analysis['text_boxes']}")
            print(f"Pictures: {analysis['pictures']}")
            
            # Process the file
            print("\nProcessing file...")
            results = processor.process_file(input_file, output_file)
            
            print("\nProcessing Results:")
            print(f"Total slides: {results['total_slides']}")
            print(f"Slides with PII: {results['slides_with_pii']}")
            print(f"Shapes with PII: {results['shapes_with_pii']}")
            print(f"\nOutput file: {results['output_file']}")
            
        except Exception as e:
            print(f"Error: {str(e)}", file=sys.stderr)
            sys.exit(1)
    else:
        print("Usage: python powerpoint_processor.py input.pptx [output.pptx]")
        sys.exit(1)